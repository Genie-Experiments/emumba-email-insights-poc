"""Slides parser.

Contains parsers for .pptx files.

"""

from pathlib import Path
from typing import Dict, List, Optional
from fsspec import AbstractFileSystem

from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document


class PptxReader(BaseReader):
    """Powerpoint parser.

    Extract text, caption images, and specify slides.

    """

    def load_data(
        self,
        file: Path,
        extra_info: Optional[Dict] = None,
        fs: Optional[AbstractFileSystem] = None,
    ) -> List[Document]:
        """Parse file."""
        from pptx import Presentation

        if fs:
            with fs.open(file) as f:
                presentation = Presentation(f)
        else:
            presentation = Presentation(file)
        result = ""
        for i, slide in enumerate(presentation.slides):
            result += f"\n\nSlide #{i}: \n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    result += f"{shape.text}\n"

        return [Document(text=result, metadata=extra_info or {})]
